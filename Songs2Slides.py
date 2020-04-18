# Import dependencies
from bs4 import BeautifulSoup
import json
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
import requests
import subprocess
import sys
import tempfile



# Gets the lyrics
def GetLyrics(artist, song):
    # Convert to lowercase
    artist = artist.lower()
    song = song.lower()

    # Replace invalid characters
    old = [" ", "!", "@", "#", "$", "%", "^", "&",     "*", "(", ")", "+", "=", "'", "?", "/", "|", "\\", ".", ",", "á", "é", "í", "ó", "ñ", "ú"]
    new = ["-", "",  "",  "",  "s", "",  "-", "-and-", "",  "",  "",  "-", "-", "",  "",  "",  "",  "",   "",  "",  "a", "e", "i", "o", "n", "u"]
    for i in range(0, len(old)):
        artist = artist.replace(old[i], new[i])
        song = song.replace(old[i], new[i])
    
    # Remove unnecessary dashes
    artist = "-".join(list(filter(lambda a: a != "", artist.split("-"))))
    song   = "-".join(list(filter(lambda a: a != "", song.split("-"))))

    # Get lyrics
    page = requests.get("https://genius.com/{0}-{1}-lyrics".format(artist, song))
    lyrics = BeautifulSoup(page.text, "html.parser").find("div", class_="lyrics").get_text()
    
    # Return lyrics
    return lyrics



# Parses the lyrics into blocks
def ParseLyrics(lyrics, settings):
    # Parse lyrics
    rawLines = lyrics.split("\n")
    lines = []
    BlockSize = settings["lines-per-slide"]
    for i in range(0, len(rawLines)):
        if (rawLines[i] == "" or rawLines[i][0] == "["):
            # Start new block on the next line
            BlockSize = settings["lines-per-slide"]
        elif (BlockSize == settings["lines-per-slide"]):
            # Start a new block
            lines.append(rawLines[i])
            BlockSize = 1
        else:
            # Continue a block
            lines[-1] = lines[-1] + "\n" + rawLines[i]
            BlockSize += 1
    return lines



# Create powerpoint
def CreatePptx(parsedLyrics, filepath, openFirst, settings):
    # Create presentation
    if (openFirst):
        try:
            prs = Presentation(filepath)
        except:
            prs = Presentation()
    else:
        prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # Set slide width and height
    prs.slide_width = Inches(settings["slide-width"])
    prs.slide_height = Inches(settings["slide-height"])
    
    # Get margins
    left = Inches(settings["margin-left"])
    top = Inches(settings["margin-top"])
    width = Inches(settings["slide-width"] - settings["margin-left"] - settings["margin-right"])
    height = Inches(settings["slide-height"] - settings["margin-top"] - settings["margin-bottom"])
    
    for lyric in parsedLyrics:
        # Add slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Apply slide formating
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(settings["slide-color"][0], settings["slide-color"][1], settings["slide-color"][2])
        
        # Add text box
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.clear()

        # Apply text formating
        tf.word_wrap = settings["word-wrap"]
        if (settings["vertical-alignment"].lower() == "top"):
            tf.vertical_anchor = MSO_ANCHOR.TOP
        elif (settings["vertical-alignment"].lower() == "bottom"):
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
        # Add pharagraph
        p = tf.paragraphs[0]
        p.text = lyric

        # Apply pharagraph formating
        p.font.name = settings["font-family"]
        p.font.size = Pt(settings["font-size"])
        p.font.bold = settings["font-bold"]
        p.font.italic = settings["font-italic"]
        p.font.color.rgb = RGBColor(settings["font-color"][0], settings["font-color"][1], settings["font-color"][2])
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = settings["line-spacing"]

    # Save powerpoint
    prs.save(filepath)



# Run CLI
if (__name__ == "__main__"):
    # Load settings
    try:
        with open(os.path.join(os.path.realpath(os.path.join(os.getcwd(), os.path.dirname(__file__))), "settings.json")) as f:
            settings = json.load(f)
    except:
        print("There was an error while loading the settings.")
        input("Press enter to exit...")
        sys.exit()
    
    # Print title
    print("Songs2Slides")
    print()

    # Get song lyrics
    lyrics = []
    song = 1
    while (True):
        # Get song information
        title = input("Song #{0} Title: ".format(song))
        artist = input("Song #{0} Artist: ".format(song))
        
        # Get song lyrics
        try:
            parsedLyrics = ParseLyrics(GetLyrics(artist, title), settings)
            if (settings["title-slides"]):
                lyrics += ["{0}\n{1}".format(title, artist)]
            lyrics += parsedLyrics
            if (lyrics[-1] != ""):
                lyrics += [""]
        except:
            print("The song could not be found. Make sure that you spelled it correctly.")
            song -= 1
        
        # Add more songs
        if (song >= 1 and input("Do you want to add another song? (y/n): ").lower() == "n"):
            break
        else:
            song += 1
    
    # Review lyrics
    if (input("Do you want to review the parsed lyrics first? (y/n): ").lower() == "y"):
        try:
            # Create temp file
            temp = tempfile.NamedTemporaryFile(mode="w+t", suffix=".txt", delete=False)
            for line in lyrics:
                temp.writelines(line)
                temp.writelines("\n\n")
            temp.close()
            
            # Open temp file and wait
            subprocess.Popen(["notepad", temp.name]).wait()

            # Read file
            with open(temp.name) as f:
                rawLines = f.read()

            # Parse lyrics
            newLyrics = []
            for song in rawLines.split("\n\n\n"):
                newLyrics += ParseLyrics(song, settings)
                if (newLyrics[-1] != ""):
                    newLyrics += [""]
            lyrics = newLyrics
        except:
            print("There was an error while reviewing the lyrics. The unrevised lyrics will be used instead.")
        finally:
            # Delete temp file
            os.remove(temp.name)

    # Get filepath
    filepath = input("Enter a filepath to save the powerpoint to: ")

    # Add extension
    if (len(filepath) == 0):
        filepath = "Untitled.pptx"
    elif (len(filepath) < 4):
        filepath += ".pptx"
    elif (len(filepath) == 4 and filepath[-4:] != ".ppt"):
        filepath += ".pptx"
    elif (len(filepath) > 4 and filepath[-5:] != ".pptx" and filepath[-4:] != ".ppt"):
        filepath += ".pptx"

    # Confirm overwrite
    if (os.path.exists(filepath)):
        openFirst = (input("This powerpoint already exists. Do you want to add on to it? (y/n): ").lower() == "y")
    else:
        openFirst = False

    # Create powerpoint
    try:
        CreatePptx(lyrics, filepath, openFirst, settings)
    except:
        print("There was an error while creating the powerpoint.")
        input("Press enter to exit...")
        sys.exit()
    
    # Open powerpoint
    if (input("Do you want to view the powerpoint now? (y/n): ").lower() == "y"):
        try:
            os.startfile(filepath)
        except:
            print("There was an error while opening the powerpoint.")
            input("Press enter to exit...")
            sys.exit()