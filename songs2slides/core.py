from dataclasses import dataclass
import pptx
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
import os
import re
import requests

@dataclass
class SongData:
    """
    Represents data about a song

    Attributes
    ----------
    title : str
        The title of the song
    artist : str
        The artist of the song
    lyrics : str
        The song's lyrics, with double newlines separating stanzas
    """

    title: str
    artist: str
    lyrics: str

def filter_lyrics(lyrics: str):
    """
    Filter raw lyrics to remove text enclosed in brackets or parenthesis

    Used by get_song_data

    Parameters
    ----------
    lyrics : str
        The raw lyrics

    Returns
    -------
    str
        The filtered lyrics
    """

    filtered = '\n' + lyrics + '\n'

    # Remove enclosed text that takes up whole numbers of lines
    filtered = re.sub(r'\n\[[^\]]*\]\n', '\n', filtered)
    filtered = re.sub(r'\n\([^\)]*\)\n', '\n', filtered)

    # Remove enclosed text that takes up partial lines
    filtered = re.sub(r'\[[^\]]*\]', '', filtered)
    filtered = re.sub(r'\([^\)]*\)', '', filtered)

    return filtered.strip()

def get_song_data(title: str, artist:str):
    """
    Get song data from an external API

    Parameters
    ----------
    title : str
        The title of the song
    artist : str
        The artist of the song

    Returns
    -------
    SongData
        The song data
    """

    # Get API URL
    url = os.getenv('API_URL')
    if url is None:
        raise Exception()
    url = url.replace('{title}', title, 1)
    url = url.replace('{artist}', artist, 1)

    # Get HTTP authorization header
    auth = os.getenv('API_AUTH', None)
    headers = { 'Authorization': auth } if auth else {}

    # Query API
    data = requests.get(url, headers=headers).json()

    # Parse response
    if 'lyrics' in data.keys():
        return SongData(data['title'], data['artist'],
                        filter_lyrics(data['lyrics']))
    else:
        raise Exception()

def parse_song_lyrics(lyrics: str, lines_per_slide: int):
    """
    Parse slide contents from the raw lyrics of a song

    Used by assemble_slides

    Parameters
    ----------
    lyrics : str
        The song lyrics
    lines_per_slide : int
        The maximum number of lines per slide

    Returns
    -------
    list of str
        The list of slide contents
    """

    slides = ['']
    line_count = 0

    for line in lyrics.strip().split('\n'):
        line = line.strip()

        if line == '':
            # Empty line represents new slide
            if line_count != 0 or len(slides) < 2 or slides[-2] != '':
                # Consecutive empty slides are not allowed
                slides += ['']
                line_count = 0

        elif lines_per_slide is None or line_count < lines_per_slide:
            # Add line to current slide
            if line_count != 0: slides[-1] += '\n'
            slides[-1] += line
            line_count += 1

        else:
            # Overflow to new slide
            slides += [line]
            line_count = 1

    # Address case where lyrics are empty
    if slides == ['', '']: slides = []

    return slides

def assemble_slides(songs: list[SongData], lines_per_slide: int = 4,
                    title_slides: bool = True, blank_slides: bool = True):
    """
    Assemble slides from a list of songs

    Paramters
    ---------
    songs : list of SongData
        The songs
    lines_per_slide : int
        The maximum number of lines per slide (default: 4)
    title_slides : bool
        Whether to include title slides before songs (default: True)
    blank_slides : bool
        Whether to include blank slides between songs (default: True)

    Returns
    -------
    list of str
        The list of slide contents
    """

    slides = []
    for song in songs:
        # Add slides for song
        if title_slides: slides += [f'{song.title}']
        slides += parse_song_lyrics(song.lyrics.upper(), lines_per_slide)
        if blank_slides: slides += ['']

    # Remove trailing blank slides
    if len(slides) and blank_slides: slides = slides[:-1]

    return slides

def create_pptx(slide_contents: list[str], filepath: str):
    """
    Create a PowerPoint from a list of slide contents

    Parameters
    ----------
    slide_contents : list of str
        The list of slide contents
    filepath : str
        The file to save the PowerPoint to
    """

    # Create presentation
    prs = pptx.Presentation()

    # Set slide aspect ratio to 16:9
    prs.slide_width = Inches(11)
    prs.slide_height = Inches(6.1875)

    # Get blank slide template
    blank_slide_layout = prs.slide_layouts[6]

    # Get textbox size parameters
    margin = Inches(0.5)
    width = prs.slide_width - 2 * margin
    height = prs.slide_height - 2 * margin

    for slide_content in slide_contents:
        # Create and format slide
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string('000000')

        # Create and format textbox
        textbox = slide.shapes.add_textbox(margin, margin, width, height)
        textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        textbox.text_frame.word_wrap = True

        # Format paragraph
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor.from_string('ffffff')
        paragraph.font.size = Pt(36)
        paragraph.alignment = PP_ALIGN.CENTER

        # Add slide content
        paragraph.text = slide_content

    # Save to file
    prs.save(filepath)
