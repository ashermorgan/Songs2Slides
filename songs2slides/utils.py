from dotenv import load_dotenv

from dataclasses import dataclass
import pptx
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
import os
import requests

@dataclass
class SongData:
    """
    Represents data about a song

    Attributes
    ----------
    title : str
        The title of the song
    artist : str
        The artist of the song
    lyrics : str
        The song's lyrics, with double newlines separating stanzas
    """

    title: str
    artist: str
    lyrics: str

def get_song_data(title: str, artist:str):
    """
    Get song data from an external API

    Parameters
    ----------
    title : str
        The title of the song
    artist : str
        The artist of the song

    Returns
    -------
    SongData
        The song data
    """

    url = os.getenv('API_URL')
    if url is None:
        raise Exception()
    url = url.replace('{title}', title, 1)
    url = url.replace('{artist}', artist, 1)
    data = requests.get(url).json()

    if 'lyrics' in data.keys():
        return SongData(data['title'], data['artist'], data['lyrics'])
    else:
        raise Exception()

def get_slide_contents(lyrics: str, lines_per_slide: int = 4):
    """
    Generate slide contents from song lyrics

    Parameters
    ----------
    lyrics : str
        The song lyrics
    lines_per_slide : int
        The maximum number of lines per slide (the default is 4)

    Returns
    -------
    list of str
        The list of slide contents
    """

    slides = ['']
    line_count = 0

    for line in lyrics.split('\n'):
        line = line.strip()

        if line == '':
            # Empty line represents new slide
            slides += ['']
            line_count = 0

        elif line_count < lines_per_slide:
            # Add line to current slide
            if line_count != 0: slides[-1] += '\n'
            slides[-1] += line
            line_count += 1

        else:
            # Overflow to new slide
            slides += [line]
            line_count = 1

    # Remove first/last slide if empty
    # len(slides) is always greater than 1 or single slide is not empty
    if slides[0] == '': slides = slides[1:]
    if slides[-1] == '': slides = slides[:-1]

    return slides

def create_pptx(slide_contents: list[str], filepath: str):
    """
    Create a PowerPoint from a list of slide contents

    Parameters
    ----------
    slide_contents : list of str
        The list of slide contents
    filepath : str
        The file to save the PowerPoint to
    """

    # Create presentation
    prs = pptx.Presentation()

    # Get blank slide template
    blank_slide_layout = prs.slide_layouts[6]

    # Get textbox size parameters
    margin = Inches(1)
    width = prs.slide_width - Inches(2)
    height = prs.slide_height - Inches(2)

    for slide_content in slide_contents:
        # Create and format slide
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string('000000')

        # Create and format textbox
        textbox = slide.shapes.add_textbox(margin, margin, width, height)
        textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        textbox.text_frame.word_wrap = True

        # Format paragraph
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor.from_string('ffffff')
        paragraph.font.size = Pt(48)
        paragraph.alignment = PP_ALIGN.CENTER

        # Add slide content
        paragraph.text = slide_content

    # Save to file
    prs.save(filepath)
