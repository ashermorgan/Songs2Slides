# Import dependencies
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
import re
import requests
from Songs2Slides import config
from unidecode import unidecode



def GetLyrics(title, artist):
    """
    Get the lyrics to a song.

    Parameters
    ----------
    title : str
        The title of the song.
    artist : str
        The name of the song's artist.
    
    Returns
    -------
    lyrics : str
        The lyrics to the song.
    title : str
        The title of the song.
    artist : str
        The name of the song's artist.
    """

    # Convert to lowercase
    artist = artist.lower()
    title = title.lower()

    # Replace invalid characters
    old = [" ", "!", "@", "#", "$", "%", "^", "&",     "*", "(", ")", "+", "=", "'", "?", "/", "|", "\\", ".", ","]
    new = ["-", "",  "",  "",  "s", "",  "-", "-and-", "",  "",  "",  "-", "-", "",  "",  "",  "",  "",   "",  ""]
    for i in range(0, len(old)):
        artist = artist.replace(old[i], new[i])
        title = title.replace(old[i], new[i])
    
    # Replace unicode characters
    artist = unidecode(artist)
    title = unidecode(title)
    
    # Remove unnecessary dashes
    artist = "-".join(list(filter(lambda a: a != "", artist.split("-"))))
    title = "-".join(list(filter(lambda a: a != "", title.split("-"))))

    # Get song info
    if (f"{artist}-{title}" in config.cachedSongs):
        # Get the cache key
        key = f"{artist}-{title}"
        
        # Get info from cache
        lyrics = config.cachedSongs[key]["lyrics"]
        title = config.cachedSongs[key]["title"]
        artist = config.cachedSongs[key]["artist"]
    else:
        # Get page from the internet
        page = requests.get(f"https://genius.com/{artist}-{title}-lyrics")
        soup = BeautifulSoup(page.text, "html.parser")
        
        # Find song info
        divs = soup.find_all("div", class_="Lyrics__Container-sc-1ynbvzw-8")
        lyrics = "\n".join([div.get_text(separator="\n") for div in divs])
        title = soup.find("h1", class_="SongHeader__Title-sc-1b7aqpg-7").get_text()
        artist = soup.find("a", class_="SongHeader__Artist-sc-1b7aqpg-9").get_text()
    
    # Return lyrics
    return lyrics, title, artist



def ParseLyrics(title, artist, settings):
    """
    Parse the lyrics of a song into slides.

    Parameters
    ----------
    title : str
        The title of the song.
    artist : str
        The name of the song's artist.
    settings : dict
        The settings to use when parsing the lyrics into slides.
    
    Returns
    -------
    list
        A list of strings containing the lyrics to the song. Each list item is one slide.
    """

    # Get lyrics
    rawLyrics, title, artist = GetLyrics(title, artist)

    # Remove content in parentheses
    if (settings["remove-parentheses"]):
        rawLyrics = re.sub(r'\s?\([^)]*\)', '', rawLyrics)
    
    # Parse Lyrics
    rawLines = rawLyrics.split("\n")

    # Add title slide
    slides = []
    if (settings["title-slides"]):
        slides += ["{0}\n{1}".format(title, artist)]

    # Parse lyrics into slides
    slideSize = settings["lines-per-slide"]
    for i in range(0, len(rawLines)):
        if (rawLines[i] == "" or rawLines[i].startswith("[")):
            # Start a new slide without content
            slides.append("")
            slideSize = 0
        elif (slideSize == settings["lines-per-slide"]):
            # Start a new slide with content
            slides.append(rawLines[i])
            slideSize = 1
        elif (slideSize == 0):
            # Continue a blank slide
            slides[-1] = slides[-1] + rawLines[i]
            slideSize += 1
        else:
            # Continue a slide
            slides[-1] = slides[-1] + "\n" + rawLines[i]
            slideSize += 1

    # Add/remove blank slide
    if (slides[-1] != "" and settings["slide-between-songs"]):
        slides += [""]
    elif (slides[-1] == "" and not settings["slide-between-songs"]):
        del slides[-1]

    # Return parsed lyrics
    return slides



def CreatePptx(parsedLyrics, filepath, settings, openFirst):
    """
    Create a PowerPoint from a list of lyrics.

    Parameters
    ----------
    parsedLyrics : list
        The list of strings containing the lyrics. Each list item will be turned into one slide.
    filepath : str
        The filepath to save the PowerPoint to.
    settings : dict
        The settings to use while creating the filepath.
    openFirst : bool
        Whether to add on to the PowerPoint file if it already exists.
    """

    if (openFirst):
        try:
            # Open presentation
            prs = Presentation(filepath)
        except:
            # Create presentation
            prs = Presentation()

            # Set slide width and height
            prs.slide_width = Inches(settings["slide-width"])
            prs.slide_height = Inches(settings["slide-height"])
    else:
        # Create presentation
        prs = Presentation()

        # Set slide width and height
        prs.slide_width = Inches(settings["slide-width"])
        prs.slide_height = Inches(settings["slide-height"])
    
    # Get blank slide
    blank_slide_layout = prs.slide_layouts[6]
    
    # Get margins
    left = Inches(settings["margin-left"])
    top = Inches(settings["margin-top"])
    width = prs.slide_width - Inches(settings["margin-left"] + settings["margin-right"])
    height = prs.slide_height - Inches(settings["margin-top"] + settings["margin-bottom"])
    
    for lyric in parsedLyrics:
        # Add slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Apply slide formating
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(settings["slide-color"][1:])
        
        # Add text box
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.clear()

        # Apply text formating
        tf.word_wrap = settings["word-wrap"]
        if (settings["vertical-alignment"].lower() == "top"):
            tf.vertical_anchor = MSO_ANCHOR.TOP
        elif (settings["vertical-alignment"].lower() == "bottom"):
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
        # Add pharagraph
        p = tf.paragraphs[0]
        p.text = lyric

        # Apply pharagraph formating
        p.font.name = settings["font-family"]
        p.font.size = Pt(settings["font-size"])
        p.font.bold = settings["font-bold"]
        p.font.italic = settings["font-italic"]
        p.font.color.rgb = RGBColor.from_string(settings["font-color"][1:])
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = settings["line-spacing"]

    # Save powerpoint
    prs.save(filepath)
