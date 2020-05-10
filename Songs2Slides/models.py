# Import dependencies
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
import re
import requests



# Gets the lyrics
def GetLyrics(title, artist):
    # Convert to lowercase
    artist = artist.lower()
    title = title.lower()

    # Replace invalid characters
    old = [" ", "!", "@", "#", "$", "%", "^", "&",     "*", "(", ")", "+", "=", "'", "?", "/", "|", "\\", ".", ",", "á", "é", "í", "ó", "ñ", "ú"]
    new = ["-", "",  "",  "",  "s", "",  "-", "-and-", "",  "",  "",  "-", "-", "",  "",  "",  "",  "",   "",  "",  "a", "e", "i", "o", "n", "u"]
    for i in range(0, len(old)):
        artist = artist.replace(old[i], new[i])
        title = title.replace(old[i], new[i])
    
    # Remove unnecessary dashes
    artist = "-".join(list(filter(lambda a: a != "", artist.split("-"))))
    title = "-".join(list(filter(lambda a: a != "", title.split("-"))))

    # Get page
    page = requests.get("https://genius.com/{0}-{1}-lyrics".format(artist, title))
    soup = BeautifulSoup(page.text, "html.parser")
    
    # Find lyrics and song information
    lyrics = soup.find("div", class_="lyrics").get_text()
    title = soup.find("h1", class_="header_with_cover_art-primary_info-title").get_text()
    artist = soup.find("a", class_="header_with_cover_art-primary_info-primary_artist").get_text()
    
    # Return lyrics
    return lyrics, title, artist



# Parses the lyrics of a song into slides
def ParseLyrics(title, artist, settings):
    # Get lyrics
    rawLyrics, title, artist = GetLyrics(title, artist)

    # Remove content in parentheses
    if (settings["remove-parentheses"]):
        rawLyrics = re.sub(r'\([^)]*\)', '', rawLyrics)
    
    # Parse Lyrics
    rawLines = rawLyrics.split("\n")

    # Remove starting and ending newlines
    del rawLines[0]
    del rawLines[0]
    del rawLines[-1]
    del rawLines[-1]

    # Add title slide
    slides = []
    if (settings["title-slides"]):
        slides += ["{0}\n{1}".format(title, artist)]

    # Parse lyrics into slides
    slideSize = settings["lines-per-slide"]
    for i in range(0, len(rawLines)):
        if (rawLines[i] == ""):
            # Start a new slide without content
            slides.append("")
            slideSize = 0
        elif (rawLines[i][0] == "["):
            # Ignore
            pass
        elif (slideSize == settings["lines-per-slide"]):
            # Start a new slide with content
            slides.append(rawLines[i])
            slideSize = 1
        elif (slideSize == 0):
            # Continue a blank slide
            slides[-1] = slides[-1] + rawLines[i]
            slideSize += 1
        else:
            # Continue a slide
            slides[-1] = slides[-1] + "\n" + rawLines[i]
            slideSize += 1

    # Add blank slide
    if (slides[-1] != ""):
        slides += [""]

    # Return parsed lyrics
    return slides



# Create powerpoint
def CreatePptx(parsedLyrics, filepath, settings, openFirst):
    if (openFirst):
        try:
            # Open presentation
            prs = Presentation(filepath)
        except:
            # Create presentation
            prs = Presentation()

            # Set slide width and height
            prs.slide_width = Inches(settings["slide-width"])
            prs.slide_height = Inches(settings["slide-height"])
    else:
        # Create presentation
        prs = Presentation()

        # Set slide width and height
        prs.slide_width = Inches(settings["slide-width"])
        prs.slide_height = Inches(settings["slide-height"])
    
    # Get blank slide
    blank_slide_layout = prs.slide_layouts[6]
    
    # Get margins
    left = Inches(settings["margin-left"])
    top = Inches(settings["margin-top"])
    width = prs.slide_width - Inches(settings["margin-left"] + settings["margin-right"])
    height = prs.slide_height - Inches(settings["margin-top"] + settings["margin-bottom"])
    
    for lyric in parsedLyrics:
        # Add slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Apply slide formating
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(settings["slide-color"][1:])
        
        # Add text box
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.clear()

        # Apply text formating
        tf.word_wrap = settings["word-wrap"]
        if (settings["vertical-alignment"].lower() == "top"):
            tf.vertical_anchor = MSO_ANCHOR.TOP
        elif (settings["vertical-alignment"].lower() == "bottom"):
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
        # Add pharagraph
        p = tf.paragraphs[0]
        p.text = lyric

        # Apply pharagraph formating
        p.font.name = settings["font-family"]
        p.font.size = Pt(settings["font-size"])
        p.font.bold = settings["font-bold"]
        p.font.italic = settings["font-italic"]
        p.font.color.rgb = RGBColor.from_string(settings["font-color"][1:])
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = settings["line-spacing"]

    # Save powerpoint
    prs.save(filepath)
